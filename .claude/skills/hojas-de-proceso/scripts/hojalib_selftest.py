# -*- coding: utf-8 -*-
"""Selftest del gate de hojas de proceso: cada criterio, en ROJO y en VERDE.

Un gate que nunca vi fallar no se si funciona. Y uno que solo vi fallar tampoco sirve: si
rechaza todo, es un cartel, no un control. Por eso cada caso va en las dos direcciones.

Arma pptx sinteticos en memoria — nada de datos de cliente, el repo es publico.

    py -3 hojalib_selftest.py          ->  codigo 0 si todos pasan
"""
import io
import os
import sys
import tempfile

from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import hoja_proceso_check as CK  # noqa: E402
import hojalib as HL  # noqa: E402

TMP = tempfile.mkdtemp(prefix="hojaselftest_")


def imagen(w_px, h_px, cuerpo_px=None, nombre=None):
    """Un PNG liso. Con `cuerpo_px`, ademas lleva la metrica de pantalla redibujada."""
    nombre = nombre or "im_%dx%d_%s.png" % (w_px, h_px, cuerpo_px)
    ruta = os.path.join(TMP, nombre)
    im = Image.new("RGB", (w_px, h_px), (200, 210, 230))
    if cuerpo_px:
        HL.guardar_pantalla(im, ruta, cuerpo_px=cuerpo_px, que_es="pantalla de prueba")
    else:
        im.save(ruta)
    return ruta


def hoja(imagenes, op="20.1", texto=None, caja=(6.0, 2.0)):
    """Un pptx de una lamina: el numero de operacion, las imagenes con su tamaño en cm, y
    opcionalmente un texto en una caja de `caja` cm."""
    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(29.7), Cm(21.0)
    s = prs.slides.add_slide(prs.slide_masters[0].slide_layouts[6])
    c = s.shapes.add_textbox(Cm(0.6), Cm(3.0), Cm(2.0), Cm(0.6))
    c.text_frame.text = op
    for ruta, (x, y, w, h) in imagenes:
        s.shapes.add_picture(ruta, Cm(x), Cm(y), Cm(w), Cm(h))
    if texto:
        t = s.shapes.add_textbox(Cm(18.0), Cm(5.0), Cm(caja[0]), Cm(caja[1]))
        tf = t.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = texto
        r.font.size = Pt(11)
    d = os.path.join(TMP, "h_%s_%d.pptx" % (op.replace(".", "_"), len(os.listdir(TMP))))
    prs.save(d)
    return d


def tipos(fallas):
    return sorted({f[2] for f in fallas})


CASOS = []


def caso(nombre, ruta, declara, espera):
    """`espera` = None para verde, o el tipo de infraccion que TIENE que aparecer."""
    CASOS.append((nombre, ruta, declara, espera))


Y0 = HL.BODY_Y + HL.BANDA_H          # arriba del bloque de imagenes
X0 = HL.IMG_X

# ── 1. jerarquia ────────────────────────────────────────────────────────────
# ROJO: tres fotos del mismo tamaño. Es el caso real de 11 de las 17 hojas de la HOTMELT.
tres_iguales = [(imagen(900, 1600), (X0 + k * 5.4, Y0, 5.0, 8.9)) for k in range(3)]
caso("3 fotos iguales, sin declarar principal", hoja(tres_iguales), {}, "sin jerarquia")
caso("3 fotos iguales, con principal declarada", hoja(tres_iguales),
     {"20.1": {"principal": 0}}, "jerarquia")

# VERDE: la principal domina, como la deja el reparto de hojalib
grande = imagen(900, 1600)
chica = imagen(900, 1600)
domina = [(grande, (X0, Y0, 5.1, 9.1)),
          (chica, (X0 + 5.4, Y0, 2.5, 4.4)),
          (chica, (X0 + 5.4, Y0 + 4.6, 2.5, 4.4))]
caso("la principal domina", hoja(domina), {"20.1": {"principal": 0}}, None)

# ROJO: declara principal la que NO es la mas grande — el caso de la 20.2 con el celular
caso("declara principal la mas chica", hoja(domina), {"20.1": {"principal": 1}}, "jerarquia")

# ── 2. legibilidad ──────────────────────────────────────────────────────────
# ROJO: pantalla de 1760 px con cuerpo 27 puesta a 5,4 cm -> 2,3 pt. Es la 20.15 del deck.
ilegible = imagen(1760, 1000, cuerpo_px=27, nombre="pant_chica.png")
caso("pantalla a 5,4 cm", hoja([(ilegible, (X0, Y0, 5.4, 3.1))]),
     {"20.1": {"principal": 0}}, "no se lee")
# VERDE: la misma pantalla, al ancho que la libreria dice que necesita (16,1 cm).
# Primero puse 16,0 "a ojo" y el caso verde salio rojo por 0,04 pt: el ancho de un caso
# verde no se elige, se pide.
_ancho_ok = HL.ancho_minimo_cm(27, 1760)
caso("la misma pantalla a %.1f cm" % _ancho_ok,
     hoja([(ilegible, (X0, Y0, _ancho_ok, _ancho_ok / 1.76))]),
     {"20.1": {"principal": 0}}, None)
# ROJO: foto SIN metrica marcada `leer` y angosta
caso("foto marcada leer, angosta",
     hoja([(imagen(1600, 900), (X0, Y0, 16.0, 9.0)),
           (imagen(1600, 900), (X0, Y0, 4.0, 2.2))]),
     {"20.1": {"principal": 0, "leer": [1]}}, "no se lee")

# ── 3. cantidad ─────────────────────────────────────────────────────────────
cuatro = [(imagen(1600, 900), (X0 + (k % 2) * 8.2, Y0 + (k // 2) * 4.6, 8.0, 4.4))
          for k in range(4)]
caso("4 imagenes", hoja(cuatro), {"20.1": {"principal": 0}}, "cantidad")
# 3 esta permitido, pero solo si hay jerarquia: tres iguales siguen siendo rojo, y por eso
# el caso verde de cantidad lleva la principal grande.
tres_con_jefe = [(imagen(1600, 900), (X0, Y0, 12.0, 6.7)),
                 (imagen(1600, 900), (X0 + 12.4, Y0, 3.4, 1.9)),
                 (imagen(1600, 900), (X0 + 12.4, Y0 + 2.2, 3.4, 1.9))]
caso("3 imagenes con jerarquia", hoja(tres_con_jefe), {"20.1": {"principal": 0}}, None)

# ── 4. texto que no entra ───────────────────────────────────────────────────
LARGO = ("Verificar la temperatura de consigna de los dos rodillos y la de proteccion, "
         "y dejar constancia en el registro de set up antes de arrancar la produccion.")
caso("texto que no entra en su caja",
     hoja([(grande, (X0, Y0, 16.0, 9.0))], texto=LARGO, caja=(4.0, 0.8)),
     {"20.1": {"principal": 0}}, "texto")
caso("el mismo texto en una caja que le da",
     hoja([(grande, (X0, Y0, 16.0, 9.0))], texto=LARGO, caja=(9.0, 4.0)),
     {"20.1": {"principal": 0}}, None)

# ── 5. una hoja sin fotos: el recuadro vacio esta PERMITIDO ──────────────────
caso("hoja sin imagenes", hoja([]), {}, None)


def main():
    ancho = max(len(c[0]) for c in CASOS)
    malos = 0
    for nombre, ruta, declara, espera in CASOS:
        fallas = CK.revisar(ruta, declara)
        t = tipos(fallas)
        if espera is None:
            ok = not fallas
            dice = "sin infracciones" if ok else "aparecio " + ", ".join(t)
        else:
            ok = espera in t
            dice = ("marca %s" % espera) if ok else "NO marco %s (dio %s)" % (espera, t or "nada")
        print("  %s  %-*s  %s" % ("ok  " if ok else "FALLA", ancho, nombre, dice))
        malos += not ok
    print("\n%d casos, %d fallan." % (len(CASOS), malos))
    return 1 if malos else 0


if __name__ == "__main__":
    sys.exit(main())
