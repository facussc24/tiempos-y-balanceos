# -*- coding: utf-8 -*-
"""Gate de las hojas de proceso: sale con codigo 1 si una hoja no cumple los criterios.

Existe por el incidente del 03/09/2026, maquina HOTMELT. Fak miro la lamina 3 y pregunto
"un celular se ve mucho mas grande que una hoja con parametros, que clase de criterio estas
aplicando". Midiendo las 17 hojas salieron DOS problemas, no uno:

  · en 11 de 17, las tres imagenes tenian exactamente el mismo tamaño: no habia jerarquia
  · **ninguna de las 7 pantallas redibujadas se leia impresa**: 2,3 a 4,5 pt contra 7 de
    minimo. Yo las habia dado por buenas mirandolas ampliadas en el monitor.

Lo que chequea, con los umbrales de `hojalib`:

  1. cada hoja declara su imagen PRINCIPAL, y esa es la mas grande
  2. lo que hay que leer, se lee: cuerpo impreso >= 7 pt
  3. como maximo 3 imagenes por hoja
  4. ningun texto se sale de su caja

    hoja_proceso_check.py <archivo.pptx> [--spec <modulo>] [--jerarquia op=idx,...]

`--spec` es un modulo python con una lista HOJAS de dicts {op, principal, leer}. Vive fuera
del repo cuando trae datos de la maquina (contraseñas de HMI, part numbers del cliente).
`--jerarquia` es el atajo para chequear sin spec: 20.2=0,20.4=0
"""
import argparse
import os
import re
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation  # noqa: E402

import hojalib as HL  # noqa: E402

EMU = 360000.0


def es_contenido(sh):
    """Logo e iconos de EPP no son imagenes de contenido: se reconocen por donde estan."""
    if sh.shape_type != 13:
        return False
    x, y = sh.left / EMU, sh.top / EMU
    if y < HL.BODY_Y - 0.4:                            # cajetin (logo)
        return False
    if x > 17.0 and sh.width / EMU < 2.5:              # banda de EPP
        return False
    return True


def texto_no_entra(sh):
    """Cuanto mas alto pide el texto que lo que la caja le da. None si entra."""
    tf = sh.text_frame
    w, h = sh.width / EMU, sh.height / EMU
    alto = 0.0
    for p in tf.paragraphs:
        texto = "".join(r.text for r in p.runs)
        if not texto.strip():
            continue
        # el run 0 de un paso es el numero "1. " y va en NEGRITA: medir el parrafo entero
        # en negrita infla el ancho y da un falso positivo.
        r0 = max(p.runs, key=lambda r: len(r.text))
        size = r0.font.size.pt if r0.font.size else 11
        ml = tf.margin_left / EMU if tf.margin_left else 0
        util = w - 2 * ml - 0.24
        if util < 0.3:
            continue
        alto += HL.lineas_wrap(texto, size, util, bool(r0.font.bold)) * 1.22 * size * HL.PT_CM
        if p.space_after:
            alto += p.space_after.pt * HL.PT_CM
        if p.space_before:
            alto += p.space_before.pt * HL.PT_CM
    mt = tf.margin_top / EMU if tf.margin_top else 0
    return alto - (h - mt + 0.02) if alto > h - mt + 0.02 else None


def revisar(ruta, declara=None):
    """Lista de infracciones (lamina, op, tipo, detalle). Vacia = pasa.

    `declara` es {op: {"principal": i, "leer": [i, ...]}}.
    """
    declara = declara or {}
    prs = Presentation(ruta)
    _, alto_bloque = HL.bloque_cm()
    bloque_cm2 = HL.IMG_W * alto_bloque
    fallas = []

    for i, s in enumerate(prs.slides):
        op = None
        for sh in s.shapes:
            if sh.has_text_frame and re.fullmatch(r"\d+\.\d+", sh.text_frame.text.strip()):
                op = sh.text_frame.text.strip()
                break

        for sh in s.shapes:                            # criterio 4, en TODAS las laminas
            if not sh.has_text_frame or not sh.text_frame.text.strip():
                continue
            sobra = texto_no_entra(sh)
            if sobra:
                fallas.append((i + 1, op or "portada", "texto",
                               "un texto pide %.2f cm mas de los que tiene la caja: %r"
                               % (sobra, sh.text_frame.text.strip()[:52])))
        if not op:
            continue

        fotos = sorted([sh for sh in s.shapes if es_contenido(sh)],
                       key=lambda q: (round(q.top / EMU, 1), round(q.left / EMU, 1)))
        if not fotos:
            continue                                   # recuadro vacio: permitido

        if len(fotos) > HL.IMAGENES_MAX:               # criterio 3
            fallas.append((i + 1, op, "cantidad",
                           "tiene %d imagenes; el maximo es %d"
                           % (len(fotos), HL.IMAGENES_MAX)))

        h = declara.get(op, {})
        leer = set(h.get("leer", []))
        for k, sh in enumerate(fotos):                 # criterio 2
            ancho = sh.width / EMU
            pt = HL.cuerpo_impreso_pt(sh.image.blob, ancho)
            if pt is not None:
                if pt < HL.CUERPO_MIN_PT - HL.TOLERANCIA_PT:
                    m = HL.metrica(sh.image.blob) or {}
                    fallas.append((i + 1, op, "no se lee",
                                   "%s: %.1f pt impreso a %.1f cm (minimo %.0f; necesita "
                                   "%.1f cm de ancho, o menos campos)"
                                   % (m.get("que_es", "una pantalla"), pt, ancho,
                                      HL.CUERPO_MIN_PT,
                                      HL.ancho_minimo_cm(m["cuerpo_px"], m["ancho_px"]))))
            elif k in leer and ancho < HL.ANCHO_MIN_LEER_CM:
                fallas.append((i + 1, op, "no se lee",
                               "la imagen %d esta marcada `leer` y mide %.1f cm (minimo %.0f)"
                               % (k + 1, ancho, HL.ANCHO_MIN_LEER_CM)))

        if "principal" not in h:                       # criterio 1
            fallas.append((i + 1, op, "sin jerarquia",
                           "la hoja no declara cual es su imagen principal"))
            continue
        idx = h["principal"]
        if not (0 <= idx < len(fotos)):
            fallas.append((i + 1, op, "sin jerarquia",
                           "declara principal=%s y la lamina tiene %d imagenes"
                           % (idx, len(fotos))))
            continue
        areas = [(x.width / EMU) * (x.height / EMU) for x in fotos]
        area_pr, tinta = areas[idx], sum(areas)
        segunda = max([a for k, a in enumerate(areas) if k != idx] or [0])
        if tinta and area_pr / tinta < HL.PRINCIPAL_MIN:
            fallas.append((i + 1, op, "jerarquia",
                           "la principal es el %.0f%% de la foto de la hoja (minimo %.0f%%)"
                           % (area_pr / tinta * 100, HL.PRINCIPAL_MIN * 100)))
        if area_pr / bloque_cm2 < HL.PRINCIPAL_MIN_BLOQUE:
            fallas.append((i + 1, op, "jerarquia",
                           "la principal ocupa %.0f%% del bloque (minimo %.0f%%): queda chica"
                           % (area_pr / bloque_cm2 * 100, HL.PRINCIPAL_MIN_BLOQUE * 100)))
        if segunda and area_pr < segunda * HL.PRINCIPAL_VENTAJA:
            fallas.append((i + 1, op, "jerarquia",
                           "la principal (%.1f cm2) no le saca %.1fx a la segunda (%.1f cm2)"
                           % (area_pr, HL.PRINCIPAL_VENTAJA, segunda)))
    return fallas


def informe(fallas):
    if not fallas:
        return "hojas de proceso: PASA. Jerarquia, legibilidad, cantidad y textos, en regla."
    out = ["HOJAS DE PROCESO — %d infraccion(es)" % len(fallas), ""]
    ancho = max(len(f[2]) for f in fallas)
    for lam, op, tipo, det in fallas:
        out.append("  lam %-2d  %-6s  %-*s  %s" % (lam, op, ancho, tipo, det))
    out += ["", "Criterios: la principal >= %.0f%% de la foto de la hoja y %.1fx la segunda "
            "· lo que hay que leer >= %.0f pt impreso · maximo %d imagenes."
            % (HL.PRINCIPAL_MIN * 100, HL.PRINCIPAL_VENTAJA, HL.CUERPO_MIN_PT,
               HL.IMAGENES_MAX)]
    return "\n".join(out)


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx")
    ap.add_argument("--spec", help="modulo con HOJAS = [{op, principal, leer}, ...]")
    ap.add_argument("--jerarquia", help="atajo sin spec, por ejemplo 20.2=0,20.4=0")
    a = ap.parse_args()

    declara = {}
    if a.spec:
        sys.path.insert(0, os.path.dirname(os.path.abspath(a.spec)) or os.getcwd())
        mod = __import__(os.path.splitext(os.path.basename(a.spec))[0])
        declara = {h["op"]: h for h in mod.HOJAS}
    if a.jerarquia:
        for par in a.jerarquia.split(","):
            op, _, idx = par.partition("=")
            declara.setdefault(op.strip(), {})["principal"] = int(idx)

    fallas = revisar(a.pptx, declara)
    print(informe(fallas))
    return 1 if fallas else 0


if __name__ == "__main__":
    sys.exit(main())
