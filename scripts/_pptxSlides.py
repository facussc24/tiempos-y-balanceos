"""
_pptxSlides.py — extender un .pptx AJENO clonando sus propias slides (python-pptx).

Por que existe: la plantilla de un cliente (VW Group, Novax) se respeta clonando slides que ya
estan en el archivo, no dibujando layouts nuevos ni eligiendo colores propios. python-pptx no
trae "duplicar slide"; la receta que funciona (informe TryOut IMG 29/08/2026, LSR Patagonia
01/09/2026) es copiar el spTree y REMAPEAR cada relacion r:embed / r:id / r:link a una relacion
nueva del slide destino. Reusar los rId del origen deja imagenes sin relacion y PowerPoint no
abre el archivo.

Funciones:
  clonar_slide(prs, idx_origen, insertar_en=None)  -> slide nuevo (al final o en `insertar_en`)
  mover_slide(prs, desde, hasta)
  quitar_shape(shape)
  set_texto(shape, texto)                          conserva el formato del primer run
  agregar_imagen(slide, ruta, left, top, max_w, max_h, align="center") -> picture, sin deformar
  agregar_texto(slide, texto, left, top, w, h, size_pt=12, bold=False, color=None, font=None,
                align="left", anchor="top")        -> textbox
  renumerar_pies(prs, prefijo="Slide Number Placeholder")
  volcar_texto(prs)                                -> [[str, ...], ...] por slide (para diffs)
  barrer(prs, patrones)                            -> [(n_slide, nombre_shape, patron, texto)]

CLI:
  python scripts/_pptxSlides.py volcar <pptx>
  python scripts/_pptxSlides.py barrer <pptx> [regex ...]   (sin regex usa la lista por defecto)

Medidas: `left/top/w/h` en PULGADAS (float) o en EMU (int). 914400 EMU = 1 pulgada.
"""
import copy
import re
import sys

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

PATRONES_DEFAULT = [
    r"#\S*(VALUE|VALOR|REF|DIV|N/A|NAME|NOMBRE)",   # errores de Excel en cualquier idioma
    r"\bTBD\b", r"\bxxx+\b", r"\bKWxx\b",
    r"LOGO DEL CLIENTE",
    r"Your presentation", r"only complete if",
    r"Lorem", r"Click to add", r"Haga clic",
]


def _emu(v):
    """float = pulgadas; int = EMU ya listo."""
    if isinstance(v, float):
        return Inches(v)
    return Emu(int(v))


# --------------------------------------------------------------------------- slides

def mover_slide(prs, desde, hasta):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    el = ids[desde]
    lst.remove(el)
    lst.insert(hasta, el)


def clonar_slide(prs, idx_origen, insertar_en=None, quitar_ole=True):
    """
    Copia la slide `idx_origen` (0-based) al final del deck y la mueve a `insertar_en` si se
    pide. Hereda layout, fondo, shapes e imagenes. NO copia notas ni animaciones. Los
    shape_id se conservan (son unicos por slide, no por deck).

    `quitar_ole=True` (default) NO copia los objetos OLE incrustados (<p:oleObj>). Probado el
    01/09/2026 con la plantilla VW Group: el objeto oculto "think-cell data - do not delete"
    clonado con su embedding compartido hace que PowerPoint NO ABRA el archivo ("could not
    open the file"); sin el, abre. Es invisible (1x1 EMU) y solo lo usa el add-in think-cell.
    """
    src = prs.slides[idx_origen]
    dst = prs.slides.add_slide(src.slide_layout)

    # los placeholders que crea el layout se van: se copian los shapes reales del origen
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)

    bg = src._element.cSld.find(qn("p:bg"))
    if bg is not None:
        dst._element.cSld.insert(0, copy.deepcopy(bg))

    tree = dst._element.cSld.spTree
    for el in src._element.cSld.spTree:
        if el.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue
        if quitar_ole and el.find(".//" + qn("p:oleObj")) is not None:
            continue
        tree.append(copy.deepcopy(el))

    # remapear TODAS las relaciones del origen a relaciones propias del destino
    rid_map = {}
    for rel in list(src.part.rels.values()):
        if rel.reltype in (RT.SLIDE_LAYOUT, RT.NOTES_SLIDE):
            continue
        if quitar_ole and rel.reltype == RT.OLE_OBJECT:
            continue   # el shape ya no se copio: sin esto quedaba una relacion huerfana al .bin (auditor 02/09)
        if rel.is_external:
            new_rid = dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new_rid = dst.part.relate_to(rel.target_part, rel.reltype)
        rid_map[rel.rId] = new_rid

    attrs = [qn("r:embed"), qn("r:id"), qn("r:link"), qn("r:pict")]
    for el in dst._element.iter():
        for a in attrs:
            v = el.get(a)
            if v is not None and v in rid_map:
                el.set(a, rid_map[v])

    if insertar_en is not None:
        mover_slide(prs, len(prs.slides) - 1, insertar_en)
    return dst


def quitar_shape(shape):
    shape._element.getparent().remove(shape._element)


def borrar_slide(prs, slide):
    """Saca la slide de la lista y suelta su relacion (si no, la parte queda huerfana)."""
    lst = prs.slides._sldIdLst
    for sldId in list(lst):
        if int(sldId.get("id")) == slide.slide_id:
            rId = sldId.get(qn("r:id"))
            lst.remove(sldId)
            prs.part.drop_rel(rId)
            return True
    return False


def reordenar(prs, slides_en_orden):
    """Deja el deck exactamente con `slides_en_orden`; las que no esten se borran."""
    quedan = {s.slide_id for s in slides_en_orden}
    for s in list(prs.slides):
        if s.slide_id not in quedan:
            borrar_slide(prs, s)
    lst = prs.slides._sldIdLst
    por_id = {int(e.get("id")): e for e in lst}
    # los ids se leen ANTES de vaciar la lista: slide_id se resuelve buscando en sldIdLst
    ids = [s.slide_id for s in slides_en_orden]
    for e in list(lst):
        lst.remove(e)
    for i in ids:
        lst.append(por_id[i])


# --------------------------------------------------------------------------- tablas

def tabla_ajustar(shape, n_filas, n_cols, anchos_in=None):
    """
    Lleva una tabla existente (clonada de la plantilla, con su estilo) a n_filas x n_cols
    copiando la ultima fila / ultima celda de cada fila. `anchos_in`: lista de anchos en
    pulgadas por columna (opcional). Devuelve el objeto table.
    """
    tbl = shape.table
    tbl_el = tbl._tbl
    grid = tbl_el.tblGrid

    def sin_ids(el):
        # PowerPoint identifica filas y columnas por <a16:rowId>/<a16:colId> dentro de un
        # <a:extLst>; una copia con el MISMO id se dibuja pero sin texto. Se sacan y listo.
        for ext in el.findall(qn("a:extLst")):
            el.remove(ext)
        return el

    # Probado el 01/09/2026: con los ids ORIGINALES conservados, las celdas nuevas de las filas
    # originales se dibujaban vacias; sacando los ids de TODAS las filas y columnas PowerPoint
    # los regenera y muestra todo. Se limpian antes de copiar.
    for gc in grid.findall(qn("a:gridCol")):
        sin_ids(gc)
    for tr in tbl_el.findall(qn("a:tr")):
        sin_ids(tr)

    cols = grid.findall(qn("a:gridCol"))
    while len(cols) < n_cols:
        grid.append(sin_ids(copy.deepcopy(cols[-1])))
        cols = grid.findall(qn("a:gridCol"))
    while len(cols) > n_cols:
        grid.remove(cols[-1])
        cols = grid.findall(qn("a:gridCol"))
    filas = tbl_el.findall(qn("a:tr"))
    for tr in filas:
        tcs = tr.findall(qn("a:tc"))
        while len(tcs) < n_cols:
            tr.append(copy.deepcopy(tcs[-1]))
            tcs = tr.findall(qn("a:tc"))
        while len(tcs) > n_cols:
            tr.remove(tcs[-1])
            tcs = tr.findall(qn("a:tc"))
    while len(filas) < n_filas:
        tbl_el.append(sin_ids(copy.deepcopy(filas[-1])))
        filas = tbl_el.findall(qn("a:tr"))
    while len(filas) > n_filas:
        tbl_el.remove(filas[-1])
        filas = tbl_el.findall(qn("a:tr"))
    if anchos_in:
        total = 0
        for gc, w in zip(grid.findall(qn("a:gridCol")), anchos_in):
            gc.set("w", str(int(Inches(w))))
            total += int(Inches(w))
        shape.width = Emu(total)
    return shape.table


def celda_texto(cell, texto, size_pt=None, bold=None, color=None):
    """
    Escribe en una celda conservando el formato del primer run. `color` (hex sin #) fija el
    color del texto: en una tabla extendida por XML las celdas nuevas pueden heredar del
    estilo un color que no se ve (01/09/2026: texto invisible en las columnas agregadas hasta
    que se fijo explicito).
    """
    tf = cell.text_frame
    p0 = tf.paragraphs[0]
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    runs = p0.runs
    if runs:
        runs[0].text = texto
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
        r = runs[0]
    else:
        r = p0.add_run()
        r.text = texto
    if size_pt is not None:
        r.font.size = Pt(size_pt)
    if bold is not None:
        r.font.bold = bold
    if color:
        r.font.color.rgb = RGBColor.from_string(color)


# --------------------------------------------------------------------------- texto

def set_texto(shape, texto):
    """
    Reemplaza el texto de un shape conservando el formato del PRIMER run del primer parrafo.
    Los parrafos y runs restantes se eliminan. Acepta '\\n' para varios parrafos: los nuevos
    copian el formato del primero.
    """
    tf = shape.text_frame
    lineas = texto.split("\n")
    p0 = tf.paragraphs[0]
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    # un campo (<a:fld>, p.ej. numero de slide) no es un run: si se deja, el texto nuevo se
    # le suma ("5" + "5" = "55"). Se saca y el texto pasa a ser fijo.
    for fld in p0._p.findall(qn("a:fld")):
        p0._p.remove(fld)
    runs = p0.runs
    if runs:
        runs[0].text = lineas[0]
        for r in runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.add_run().text = lineas[0]
    for linea in lineas[1:]:
        p_new = copy.deepcopy(p0._p)
        tf._txBody.append(p_new)
        p_obj = tf.paragraphs[-1]
        p_obj.runs[0].text = linea
        for r in p_obj.runs[1:]:
            r._r.getparent().remove(r._r)


def agregar_texto(slide, texto, left, top, w, h, size_pt=12, bold=False, color=None,
                  font=None, align="left", anchor="top", wrap=True):
    tb = slide.shapes.add_textbox(_emu(left), _emu(top), _emu(w), _emu(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                          "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    for margen in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, margen, Inches(0.03))
    lineas = texto.split("\n")
    for i, linea in enumerate(lineas):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
        r = p.add_run()
        r.text = linea
        r.font.size = Pt(size_pt)
        r.font.bold = bold
        if font:
            r.font.name = font
        if color:
            r.font.color.rgb = RGBColor.from_string(color)
    return tb


# --------------------------------------------------------------------------- imagenes

def agregar_imagen(slide, ruta, left, top, max_w, max_h, align="center", valign="middle"):
    """Inserta la imagen ajustada DENTRO de la caja (left, top, max_w, max_h) sin deformarla."""
    left, top, max_w, max_h = (_emu(left), _emu(top), _emu(max_w), _emu(max_h))
    with Image.open(ruta) as im:
        iw, ih = im.size
    escala = min(max_w / iw, max_h / ih)
    w, h = int(iw * escala), int(ih * escala)
    x = {"left": left, "center": left + (max_w - w) // 2, "right": left + max_w - w}[align]
    y = {"top": top, "middle": top + (max_h - h) // 2, "bottom": top + max_h - h}[valign]
    return slide.shapes.add_picture(ruta, Emu(x), Emu(y), Emu(w), Emu(h))


# --------------------------------------------------------------------------- pies y control

def renumerar_pies(prs, prefijo="Slide Number Placeholder"):
    """
    En plantillas donde el numero de pagina es TEXTO fijo (no campo), lo reescribe como
    indice+1 en cada slide que tenga el shape. Si el shape ya tiene un campo <a:fld
    type="slidenum"> PowerPoint lo actualiza solo y aca NO se toca (la plantilla VW Group
    trae campo: renumerarla a mano duplicaba el digito). Devuelve cuantos toco.
    """
    n = 0
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if not (sh.name.startswith(prefijo) and sh.has_text_frame):
                continue
            if sh._element.findall(".//" + qn("a:fld")):
                continue
            if sh.text_frame.text.strip().isdigit() or sh.text_frame.text.strip() == "":
                set_texto(sh, str(i + 1))
                n += 1
    return n


def es_campo_slidenum(shape):
    return bool(shape.has_text_frame and shape._element.findall(".//" + qn("a:fld")))


def _textos_shape(sh):
    if sh.has_text_frame:
        t = sh.text_frame.text.strip()
        return [t] if t else []
    if getattr(sh, "has_table", False) and sh.has_table:
        out = []
        for r in sh.table.rows:
            fila = " | ".join(c.text.strip().replace("\n", " / ") for c in r.cells)
            if fila.replace("|", "").strip():
                out.append(fila)
        return out
    if sh.shape_type is not None and hasattr(sh, "shapes"):  # grupo
        out = []
        for s in sh.shapes:
            out.extend(_textos_shape(s))
        return out
    return []


def volcar_texto(prs):
    """Lista por slide de 'nombre_shape :: texto', para comparar antes/despues."""
    out = []
    for slide in prs.slides:
        lineas = []
        for sh in slide.shapes:
            for t in _textos_shape(sh):
                lineas.append(f"{sh.name} :: {t}")
        out.append(lineas)
    return out


def barrer(prs, patrones=None):
    patrones = patrones or PATRONES_DEFAULT
    regs = [re.compile(p, re.IGNORECASE) for p in patrones]
    hallazgos = []
    for i, slide in enumerate(prs.slides, 1):
        for sh in slide.shapes:
            for t in _textos_shape(sh):
                for rg in regs:
                    if rg.search(t):
                        hallazgos.append((i, sh.name, rg.pattern, t[:120]))
    return hallazgos


def _cli():
    if len(sys.argv) < 3:
        sys.exit(__doc__)
    cmd, ruta = sys.argv[1], sys.argv[2]
    prs = Presentation(ruta)
    if cmd == "volcar":
        for i, lineas in enumerate(volcar_texto(prs), 1):
            print(f"===== SLIDE {i}")
            for l in lineas:
                print("   ", l)
    elif cmd == "barrer":
        pats = sys.argv[3:] or None
        h = barrer(prs, pats)
        for n, nombre, pat, t in h:
            print(f"slide {n:>2}  {nombre!r:36} /{pat}/  ->  {t}")
        print(f"\n{len(h)} hallazgo(s)")
        sys.exit(1 if h else 0)
    else:
        sys.exit(__doc__)


if __name__ == "__main__":
    _cli()
